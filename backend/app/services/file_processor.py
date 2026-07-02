import os
import uuid
import subprocess
import threading
from typing import List, Dict, Any
from PIL import Image
import numpy as np
from config.config import settings

# Whisper 模型缓存（避免每次视频处理都重新加载）
_whisper_model = None
_whisper_model_name = None
_whisper_model_lock = threading.Lock()


def _get_whisper_model(model_name="small"):
    """懒加载 Whisper 模型，全局单例缓存"""
    global _whisper_model, _whisper_model_name
    if _whisper_model is not None and _whisper_model_name == model_name:
        return _whisper_model
    with _whisper_model_lock:
        if _whisper_model is not None and _whisper_model_name == model_name:
            return _whisper_model
        try:
            import whisper
            if hasattr(whisper, 'load_model'):
                print(f"[INFO] Loading Whisper model '{model_name}'...")
                _whisper_model = whisper.load_model(model_name)
                _whisper_model_name = model_name
                print(f"[INFO] Whisper model '{model_name}' loaded successfully")
                return _whisper_model
        except Exception as e:
            print(f"[ERROR] Failed to load Whisper model '{model_name}': {e}")
    return None


class FileProcessor:
    """文件处理器类"""
    
    @staticmethod
    def process_document(file_path: str, file_type: str) -> List[str]:
        """处理文档文件"""
        chunks = []
        try:
            if file_type == "pdf":
                # 使用pdfplumber解析PDF
                import pdfplumber
                with pdfplumber.open(file_path) as pdf:
                    text = ""
                    for page in pdf.pages:
                        text += page.extract_text() or ""
                chunks = FileProcessor._split_text(text)
            elif file_type in ["docx", "doc"]:
                # 使用python-docx解析Word
                from docx import Document
                doc = Document(file_path)
                text = "\n".join([para.text for para in doc.paragraphs])
                chunks = FileProcessor._split_text(text)
            elif file_type in ["xlsx", "xls", "ods"]:
                # 使用openpyxl解析Excel和ODS
                from openpyxl import load_workbook
                wb = load_workbook(file_path)
                text = ""
                for sheet in wb.sheetnames:
                    ws = wb[sheet]
                    for row in ws.iter_rows(values_only=True):
                        row_text = "\t".join([str(cell) if cell is not None else "" for cell in row])
                        text += row_text + "\n"
                chunks = FileProcessor._split_text(text)
            elif file_type == "csv":
                # 处理CSV文件
                import csv
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    reader = csv.reader(f)
                    text = "\n".join(["\t".join(row) for row in reader])
                chunks = FileProcessor._split_text(text)
            elif file_type == "rtf":
                # 使用striprtf解析RTF
                try:
                    from striprtf.striprtf import rtf_to_text
                    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                        rtf_content = f.read()
                    text = rtf_to_text(rtf_content)
                    chunks = FileProcessor._split_text(text)
                except ImportError:
                    # 如果striprtf未安装，尝试使用unrtf命令行工具
                    try:
                        result = subprocess.run(["unrtf", file_path], capture_output=True, text=True)
                        text = result.stdout
                        chunks = FileProcessor._split_text(text)
                    except:
                        print("无法解析RTF文件")
            elif file_type in ["ppt", "pptx", "pps", "odp"]:
                # 使用python-pptx解析PPT
                try:
                    from pptx import Presentation
                    prs = Presentation(file_path)
                    text = ""
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text += shape.text + "\n"
                    chunks = FileProcessor._split_text(text)
                except Exception as e:
                    print(f"解析PPT失败: {e}")
            elif file_type == "odt":
                # 使用odfpy解析ODT
                try:
                    from odf import text, teletype
                    from odf.opendocument import load
                    doc = load(file_path)
                    paragraphs = doc.getElementsByType(text.P)
                    text_content = "\n".join([teletype.extractText(p) for p in paragraphs])
                    chunks = FileProcessor._split_text(text_content)
                except ImportError:
                    print("odfpy未安装，无法解析ODT文件")
            elif file_type in ["md", "txt"]:
                # 直接读取文本文件
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    text = f.read()
                chunks = FileProcessor._split_text(text)
        except Exception as e:
            print(f"处理文档失败: {e}")
        return chunks
    
    @staticmethod
    def process_image(file_path: str) -> Dict[str, Any]:
        """处理图片文件"""
        result = {
            "text": "",
            "faces": []
        }
        try:
            # OCR识别文字（使用原始分辨率，Tesseract 需要足够细节）
            if settings.ENABLE_OCR:
                try:
                    import pytesseract
                    img = Image.open(file_path)
                    result["text"] = pytesseract.image_to_string(img)
                except Exception as e:
                    print(f"OCR识别失败: {e}")

            # 提取人脸特征（可配置，使用原始文件路径）
            if settings.ENABLE_FACE_DETECTION:
                try:
                    from app.services.face_service import FaceService
                    face_service = FaceService()

                    if hasattr(face_service, 'insightface_available') and face_service.insightface_available:
                        result["faces"] = face_service.extract_face_features(file_path)
                    else:
                        result["faces"] = []
                except Exception as e:
                    result["faces"] = []
            else:
                result["faces"] = []
        except Exception as e:
            pass
        return result
    
    @staticmethod
    def process_video(file_path: str) -> Dict[str, Any]:
        """处理视频文件"""
        result = {
            "frames": [],
            "audio_text": ""
        }
        try:
            # 视频抽帧
            frames_dir = os.path.join(settings.TEMP_DIR, f"frames_{str(uuid.uuid4())}")
            os.makedirs(frames_dir, exist_ok=True)
            
            # 使用ffmpeg抽帧
            cmd = [
                "ffmpeg", "-i", file_path,
                "-vf", f"fps=1/{settings.VIDEO_FRAME_INTERVAL}",
                "-vframes", str(settings.MAX_FRAMES_PER_VIDEO),
                f"{frames_dir}/frame_%04d.jpg"
            ]
            subprocess.run(cmd, capture_output=True, text=True)
            
            # 处理帧图片
            frame_files = sorted([f for f in os.listdir(frames_dir) if f.endswith('.jpg')])
            for i, frame_file in enumerate(frame_files):
                frame_path = os.path.join(frames_dir, frame_file)
                frame_result = FileProcessor.process_image(frame_path)
                result["frames"].append({
                    "path": frame_path,
                    "timestamp": i * settings.VIDEO_FRAME_INTERVAL,
                    "text": frame_result["text"],
            
                })
            
            # 提取音频并转写（可配置）
            if settings.ENABLE_VIDEO_AUDIO:
                audio_path = os.path.join(settings.TEMP_DIR, f"audio_{str(uuid.uuid4())}.mp3")
                try:
                    cmd = ["ffmpeg", "-i", file_path, "-q:a", "0", "-map", "a", audio_path]
                    subprocess.run(cmd, capture_output=True, text=True)
                    
                    # 使用whisper转写音频（使用全局缓存模型，避免重复加载）
                    try:
                        model = _get_whisper_model("small")
                        if model:
                            transcript = model.transcribe(audio_path)
                            result["audio_text"] = transcript["text"]
                        
                    except Exception as e:
                        result["audio_text"] = ""
                except Exception as e:
                    result["audio_text"] = ""
                
                # 清理音频文件
                if os.path.exists(audio_path):
                    os.remove(audio_path)
            else:
                result["audio_text"] = ""
            
            # 清理临时文件
            FileProcessor._clean_temp_dir(frames_dir)
        except Exception as e:
            print(f"处理视频失败: {e}")
        return result
    
    @staticmethod
    def _split_text(text: str) -> List[str]:
        """文本分块"""
        chunks = []
        words = text.split()
        current_chunk = []
        current_length = 0
        
        for word in words:
            current_chunk.append(word)
            current_length += len(word) + 1  # +1 for space
            
            if current_length >= settings.CHUNK_SIZE:
                chunks.append(" ".join(current_chunk))
                # 保留重叠部分
                overlap_size = min(len(current_chunk), settings.CHUNK_OVERLAP // 10)
                current_chunk = current_chunk[-overlap_size:]
                current_length = sum(len(w) + 1 for w in current_chunk)
        
        if current_chunk:
            chunks.append(" ".join(current_chunk))
        
        return chunks
    
    @staticmethod
    def _clean_temp_dir(directory: str):
        """清理临时目录"""
        if os.path.exists(directory):
            for file in os.listdir(directory):
                file_path = os.path.join(directory, file)
                if os.path.isfile(file_path):
                    os.remove(file_path)
            os.rmdir(directory)
